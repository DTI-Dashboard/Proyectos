import traceback as _tb
import sys as _sys

_error_file = "sync_last_error.txt"

try:
    import os
    import json
    import urllib.request
    import re
    import csv
    from datetime import datetime, timezone

    TOKEN = os.environ["SMARTSHEET_TOKEN"]

    def fmt_fecha(f):
        if not f: return ""
        f = str(f).split("T")[0]
        for fmt in ("%Y-%m-%d", "%m/%d/%Y", "%d/%m/%Y"):
            try: return datetime.strptime(f, fmt).strftime("%d/%m/%Y")
            except ValueError: continue
        return str(f)

    def calcular_semaforo(real, esp, fin):
        hoy = datetime.today()
        try:
            fin_vencido = datetime.strptime(fmt_fecha(fin), "%d/%m/%Y") < hoy
        except Exception:
            fin_vencido = False
        if real >= 100:       return "Terminado"
        elif fin_vencido:     return "Atrasado"
        elif real >= esp:     return "A tiempo"
        elif real <= esp - 11: return "Atrasado"
        else:                 return "En riesgo"

    def fetch_report(report_id, modulo_label, filter_mode="sheet_match"):
        """
        filter_mode:
          'sheet_match' → fila es proyecto si sheet_name está dentro de primary (IMC)
          'parent_row'  → fila es proyecto si parentId es None (DDI/Desarrollo)
        """
        url = f"https://api.smartsheet.com/2.0/reports/{report_id}?pageSize=10000"
        req = urllib.request.Request(url, headers={
            "Authorization": f"Bearer {TOKEN}",
            "Accept": "application/json"
        })
        with urllib.request.urlopen(req) as resp:
            data = json.loads(resp.read())

        col_map = {}
        for c in data.get("columns", []):
            col_id = c.get("virtualId") or c.get("virtualColumnId") or c.get("id")
            if col_id:
                col_map[str(col_id)] = c.get("title", "")

        proyectos = []
        seen_sheets = set()
        hoy = datetime.today()

        for row in data.get("rows", []):
            cells = {}
            for cell in row.get("cells", []):
                col_id = str(cell.get("virtualColumnId") or cell.get("columnId") or "")
                col_title = col_map.get(col_id, "")
                val = cell.get("displayValue") or cell.get("value") or ""
                if col_title:
                    cells[col_title] = val

            sheet_name = str(cells.get("Sheet Name", "") or cells.get("Nombre de la hoja", "")).strip()
            primary    = str(cells.get("Primary", "")).strip()
            avance_raw = cells.get("Avance", "")
            inicio     = cells.get("Fecha de inicio", "")
            fin        = cells.get("Fecha de finalización", "")
            parent_id  = row.get("parentId")

            if not primary:
                continue

            if filter_mode == "sheet_match":
                # IMC: el nombre del sheet está contenido en el Primary
                if not sheet_name or sheet_name.lower() not in primary.lower():
                    continue
            elif filter_mode == "first_per_sheet":
                # DDI: tomar solo la primera fila de cada hoja (= el proyecto resumen)
                if not sheet_name:
                    continue
                if sheet_name in seen_sheets:
                    continue
                seen_sheets.add(sheet_name)

            try:
                real = int(float(str(avance_raw).replace("%", "").strip()))
            except (ValueError, TypeError):
                real = 0

            try:
                d_ini = datetime.strptime(fmt_fecha(inicio), "%d/%m/%Y")
                d_fin = datetime.strptime(fmt_fecha(fin), "%d/%m/%Y")
                total = (d_fin - d_ini).days
                trans = (hoy - d_ini).days
                esp   = max(0, min(100, int((trans / total) * 100))) if total > 0 else 0
            except Exception:
                esp = real

            sem = calcular_semaforo(real, esp, fin)

            proyectos.append({
                "nombre":      primary,
                "inicio":      fmt_fecha(inicio),
                "fin":         fmt_fecha(fin),
                "real":        real,
                "esp":         esp,
                "sem":         sem,
                "area":        "",
                "responsable": "",
                "solicitante": "",
                "descripcion": "",
                "url_plan": ""
            })

        print(f"✅ {len(proyectos)} proyectos {modulo_label} de Smartsheet")
        return proyectos

    def load_csv(filepath):
        info = {}
        if os.path.exists(filepath):
            with open(filepath, "r", encoding="utf-8-sig") as f:
                reader = csv.DictReader(f)
                for row in reader:
                    nombre = row.get("nombre", "").strip()
                    if nombre:
                        info[nombre.lower()] = {
                            "area":        row.get("area", ""),
                            "responsable": row.get("responsable", ""),
                            "solicitante": row.get("solicitante", ""),
                            "descripcion": row.get("descripcion", ""),
                            "url_plan":    row.get("url_plan", ""),
                        }
            print(f"✅ {filepath}: {len(info)} entradas")
        else:
            print(f"ℹ️  {filepath} no existe aún")
        return info

    def merge_info(proyectos, info_extra):
        for p in proyectos:
            extra = info_extra.get(p["nombre"].lower(), {})
            for key in ("area", "responsable", "solicitante", "descripcion", "url_plan"):
                if extra.get(key): p[key] = extra[key]

    def proyectos_to_js(proyectos, var_name):
        js_rows = []
        for p in proyectos:
            def esc(s): return s.replace("\\","\\\\").replace("'","\\'").replace('"','\\"')
            desc = p["descripcion"] if p["descripcion"] else "Sincronizado desde Smartsheet"
            js_rows.append(
                f'  {{tipo:"Proyecto",nombre:"{esc(p["nombre"])}",area:"{esc(p["area"])}",'
                f'inicio:"{p["inicio"]}",fin:"{p["fin"]}",real:{p["real"]},'
                f'esp:{p["esp"]},sem:"{p["sem"]}",responsable:"{esc(p["responsable"])}",'
                f'solicitante:"{esc(p["solicitante"])}",descripcion:"{esc(desc)}",'
                f'url_plan:"{esc(p["url_plan"])}"}}'
            )
        return f"const {var_name} = [\n" + ",\n".join(js_rows) + "\n];"

    # ── Configuración de módulos ──────────────────────────────────
    MODULOS = [
        {
            "report_id":   "1111996261945220",
            "label":       "IMC",
            "var":         "IMC_DATA",
            "csv":         "info_proyectos.csv",
            "filter_mode": "sheet_match"
        },
        {
            "report_id":   "2916326996660100",
            "label":       "Desarrollo",
            "var":         "DEV_DATA",
            "csv":         "info_dev.csv",
            "filter_mode": "first_per_sheet"
        },
        {
            "report_id":   "6503854223871876",
            "label":       "IA",
            "var":         "IA_DATA",
            "csv":         "info_ia.csv",
            "filter_mode": "first_per_sheet"
        },
    ]

    # ── Leer index.html ───────────────────────────────────────────
    with open("index.html", "r", encoding="utf-8") as f:
        html = f.read()

    # ── Procesar cada módulo ──────────────────────────────────────
    for mod in MODULOS:
        try:
            proyectos  = fetch_report(mod["report_id"], mod["label"], mod["filter_mode"])
            info_extra = load_csv(mod["csv"])
            merge_info(proyectos, info_extra)
            new_data   = proyectos_to_js(proyectos, mod["var"])
            pattern    = rf"const {mod['var']}\s*=\s*\[.*?\];"
            match = re.search(pattern, html, flags=re.DOTALL)
            if not match:
                print(f"⚠️  Patrón no encontrado para {mod['var']} — saltando")
                continue
            html = html[:match.start()] + new_data + html[match.end():]
            print(f"✅ {mod['var']} inyectado: {len(proyectos)} proyectos")
        except Exception as e:
            import traceback
            print(f"❌ Error procesando {mod['label']}: {e}")
            traceback.print_exc()
            raise

    # ── Actividades desde CSV ────────────────────────────────────────
    act_csv = "actividades.csv"
    act_data = []
    if os.path.exists(act_csv):
        with open(act_csv, "r", encoding="utf-8-sig") as f:
            reader = csv.DictReader(f)
            for row in reader:
                act_data.append({
                    "id":        row.get("id",""),
                    "nombre":    row.get("nombre",""),
                    "actividad": row.get("actividad",""),
                    "area":      row.get("area",""),
                    "prioridad": row.get("prioridad",""),
                    "fechaFin":  row.get("fechaFin",""),
                    "semana":    row.get("semana",""),
                    "estatus":   row.get("estatus",""),
                    "notas":     row.get("notas",""),
                    "createdAt": row.get("createdAt",""),
                    "updatedAt": row.get("updatedAt",""),
                })
        print(f"✅ actividades.csv: {len(act_data)} registros")
    else:
        print("ℹ️  actividades.csv no existe aún")

    def esc_js(s):
        return (s or "").replace("\\","\\\\").replace("'","\\'").replace('"','\\"')

    act_js_rows = []
    for a in act_data:
        act_js_rows.append(
            f'  {{id:"{esc_js(a["id"])}",nombre:"{esc_js(a["nombre"])}",actividad:"{esc_js(a["actividad"])}",'
            f'area:"{esc_js(a["area"])}",prioridad:"{esc_js(a["prioridad"])}",fechaFin:"{esc_js(a["fechaFin"])}",'
            f'semana:"{esc_js(a["semana"])}",estatus:"{esc_js(a["estatus"])}",notas:"{esc_js(a["notas"])}",'
            f'createdAt:"{esc_js(a["createdAt"])}",updatedAt:"{esc_js(a["updatedAt"])}"}}'
        )
    act_js = "const ACT_DATA = [\n" + ",\n".join(act_js_rows) + "\n]; // INYECTADO POR SYNC.PY"
    pattern_act = r"const ACT_DATA = \[.*?\]; // INYECTADO POR SYNC\.PY"
    match_act = re.search(pattern_act, html, flags=re.DOTALL)
    if match_act:
        html = html[:match_act.start()] + act_js + html[match_act.end():]
        print(f"✅ ACT_DATA inyectado: {len(act_data)} actividades")
    else:
        print("⚠️  Patrón ACT_DATA no encontrado en HTML")


    # ── Generar cierre_base.pptx con datos IMC actuales ──────────────────────
    try:
        from pptx import Presentation
        from pptx.util import Inches as I, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        import lxml.etree as etree, io, base64, urllib.request as _ur

        # Leer datos IMC del HTML ya actualizado
        import re as _re
        idx_imc = html.find('const IMC_DATA = [')
        depth = 0
        i_pos = idx_imc + len('const IMC_DATA = ')
        for j in range(i_pos, i_pos+500000):
            if html[j] == '[': depth += 1
            elif html[j] == ']':
                depth -= 1
                if depth == 0:
                    block = html[i_pos:j+1]
                    break
        matches = _re.findall(
            r'\{tipo:"([^"]*)",nombre:"([^"]*)",area:"([^"]*)",inicio:"([^"]*)",fin:"([^"]*)",real:(\d+),esp:(\d+),sem:"([^"]*)"',
            block
        )
        datos = [{'tipo':m[0],'nombre':m[1],'area':m[2],'inicio':m[3],'fin':m[4],'real':int(m[5]),'esp':int(m[6]),'sem':m[7]} for m in matches]

        # Descargar template
        tpl_url = 'https://raw.githubusercontent.com/DTI-Dashboard/Proyectos/main/template_cierre.pptx'
        with _ur.urlopen(tpl_url) as _r:
            tpl_bytes = _r.read()

        from datetime import date, timedelta
        hoy_d = date.today()
        dow = hoy_d.weekday()
        dias_atras = (dow - 3) % 7
        jP = hoy_d - timedelta(days=dias_atras)
        jA = jP + timedelta(days=7)
        MESES = ['Enero','Febrero','Marzo','Abril','Mayo','Junio','Julio','Agosto','Septiembre','Octubre','Noviembre','Diciembre']
        rango = f"{jP.day} de {MESES[jP.month-1]} al {jA.day} de {MESES[jA.month-1]}"

        prs = Presentation(io.BytesIO(tpl_bytes))
        NAVY=RGBColor(0x1A,0x23,0x40); GOLD=RGBColor(0xC9,0xA8,0x4C)
        WHITE=RGBColor(0xFF,0xFF,0xFF); GRAY=RGBColor(0x55,0x55,0x55)
        sem_col = {'Terminado':RGBColor(0x22,0xC5,0x5E),'A tiempo':RGBColor(0x22,0xC5,0x5E),'Atrasado':RGBColor(0xEF,0x44,0x44),'Stand By':RGBColor(0x11,0x11,0x11),'Riesgo de atraso':RGBColor(0xF5,0x9E,0x0B)}
        sem_hex = {'Terminado':'22C55E','A tiempo':'22C55E','Atrasado':'EF4444','Stand By':'111111','Riesgo de atraso':'F59E0B'}
        sem_txt = {'Terminado':'Terminado','A tiempo':'A tiempo','Atrasado':'Atrasado','Stand By':'Stand by','Riesgo de atraso':'Riesgo de atraso'}

        def sc(p):
            s=(p['sem'] or '').lower()
            if 'terminado' in s: return 'Terminado'
            if 'tiempo' in s: return 'A tiempo'
            if 'riesgo' in s: return 'Riesgo de atraso'
            if 'stand' in s: return 'Stand By'
            return 'Atrasado'

        def cell_txt(cell, text, sz=9, bold=False, color=None, align=PP_ALIGN.LEFT):
            tf=cell.text_frame; tf.word_wrap=True; tf.clear()
            p=tf.paragraphs[0]; p.alignment=align
            r=p.add_run(); r.text=str(text)
            r.font.size=Pt(sz); r.font.bold=bold; r.font.name='Calibri'
            if color: r.font.color.rgb=color

        def bar_fill(cell, real, sc_str):
            tc=cell._tc
            NS='http://schemas.openxmlformats.org/drawingml/2006/main'
            tcPr=tc.find(f'{{{NS}}}tcPr')
            if tcPr is None:
                tcPr=etree.Element(f'{{{NS}}}tcPr'); tc.insert(0,tcPr)
            for f in list(tcPr):
                if 'Fill' in f.tag or 'noFill' in f.tag: tcPr.remove(f)
            pct=min(float(real or 0),100)
            h=sem_hex.get(sc_str,'888888')
            if pct>=100:
                fill=etree.SubElement(tcPr,f'{{{NS}}}solidFill')
                etree.SubElement(fill,f'{{{NS}}}srgbClr').set('val',h)
            else:
                stop=int(pct*1000); s3=min(stop+500,100000)
                gf=etree.SubElement(tcPr,f'{{{NS}}}gradFill')
                gl=etree.SubElement(gf,f'{{{NS}}}gsLst')
                for pos,col in [(0,h),(stop,h),(s3,'E2E2E2'),(100000,'E2E2E2')]:
                    g=etree.SubElement(gl,f'{{{NS}}}gs'); g.set('pos',str(pos))
                    etree.SubElement(g,f'{{{NS}}}srgbClr').set('val',col)
                lin=etree.SubElement(gf,f'{{{NS}}}lin'); lin.set('ang','0'); lin.set('scaled','0')

        # SLIDE 1
        s1=prs.slides[0]
        for shape in s1.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if 'semanal del' in run.text:
                            run.text=f'Cierre semanal del {rango}'
                            shape.width=I(7.5); shape.left=I(1.7)

        # SLIDE 2
        s2=prs.slides[1]
        for shape in list(s2.shapes):
            if hasattr(shape,'image'):
                shape._element.getparent().remove(shape._element)

        tot=len(datos); ter=sum(1 for p in datos if sc(p)=='Terminado')
        ati=sum(1 for p in datos if sc(p)=='A tiempo'); rie=sum(1 for p in datos if sc(p)=='Riesgo de atraso')
        atr=sum(1 for p in datos if sc(p)=='Atrasado')
        pcr=sum(1 for p in datos if p['real']>=90 and sc(p)!='Terminado')
        stb=sum(1 for p in datos if sc(p)=='Stand By')
        kpis=[('Total de proyectos',tot),('Terminado',ter),('A tiempo',ati),('En riesgo\nde atraso',rie),('Atrasados',atr),('Por cerrar',pcr),('Stand By',stb)]
        kw=1.345; kh=0.62; ky=0.62; kg=0.055; ksx=0.08
        for ii,(lbl,val) in enumerate(kpis):
            x=ksx+ii*(kw+kg)
            bx=s2.shapes.add_shape(1,I(x),I(ky),I(kw),I(kh))
            bx.fill.solid(); bx.fill.fore_color.rgb=WHITE
            bx.line.color.rgb=GOLD; bx.line.width=Pt(1.2)
            vt=s2.shapes.add_textbox(I(x),I(ky+0.03),I(kw),I(0.32))
            vp=vt.text_frame.paragraphs[0]; vp.alignment=PP_ALIGN.CENTER
            vr=vp.add_run(); vr.text=str(val); vr.font.size=Pt(20); vr.font.bold=True; vr.font.name='Calibri'; vr.font.color.rgb=NAVY
            lt=s2.shapes.add_textbox(I(x),I(ky+0.37),I(kw),I(0.24))
            lt.text_frame.word_wrap=True
            lp=lt.text_frame.paragraphs[0]; lp.alignment=PP_ALIGN.CENTER
            lr=lp.add_run(); lr.text=lbl; lr.font.size=Pt(6.5); lr.font.name='Calibri'; lr.font.color.rgb=GRAY
        tx=s2.shapes.add_textbox(I(1.75),I(0.04),I(6.5),I(0.34))
        p_tx=tx.text_frame.paragraphs[0]; p_tx.alignment=PP_ALIGN.CENTER
        r_tx=p_tx.add_run(); r_tx.text='ÁREA DE IMPLEMENTACIÓN'
        r_tx.font.size=Pt(15); r_tx.font.bold=True; r_tx.font.name='Calibri'; r_tx.font.color.rgb=NAVY
        tx2=s2.shapes.add_textbox(I(0),I(0.38),I(10),I(0.18))
        p2=tx2.text_frame.paragraphs[0]; p2.alignment=PP_ALIGN.CENTER
        r2=p2.add_run(); r2.text=f'Cierre semanal del {rango}'
        r2.font.size=Pt(8.5); r2.font.name='Calibri'; r2.font.color.rgb=GRAY
        colW=[I(0.83),I(2.70),I(0.73),I(0.73),I(0.58),I(0.65),I(2.42),I(0.92)]
        TBL_X=0.08; TBL_Y=1.30; HDR_H=0.22; ROW_H=0.215
        tblW=sum(colW); nrows=len(datos)+1
        tblS=s2.shapes.add_table(nrows,8,I(TBL_X),I(TBL_Y),tblW,I(HDR_H+len(datos)*ROW_H))
        tbl=tblS.table
        for ci,w in enumerate(colW): tbl.columns[ci].width=w
        tbl.rows[0].height=I(HDR_H)
        hdrs=['TIPO','PROYECTO','INICIO','FIN','% REAL','% ESPERADO','BARRA DE AVANCE','SEMÁFORO']
        for ci,h in enumerate(hdrs):
            cell=tbl.cell(0,ci)
            cell_txt(cell,h,sz=7,bold=True,color=NAVY,align=PP_ALIGN.CENTER)
            cell.fill.solid(); cell.fill.fore_color.rgb=GOLD
        for ri,p in enumerate(datos):
            tbl.rows[ri+1].height=I(ROW_H)
            s=sc(p); sc_c=sem_col.get(s,RGBColor(0x88,0x88,0x88))
            bg=RGBColor(0xF5,0xF7,0xFA) if ri%2==0 else WHITE
            vals=[p['tipo'],p['nombre'],p['inicio'],p['fin'],f"{p['real']}%",f"{p['esp']}%",f"{p['real']}%",sem_txt.get(s,s)]
            aligns=[PP_ALIGN.CENTER,PP_ALIGN.LEFT,PP_ALIGN.CENTER,PP_ALIGN.CENTER,PP_ALIGN.CENTER,PP_ALIGN.CENTER,PP_ALIGN.RIGHT,PP_ALIGN.CENTER]
            colors=[NAVY,NAVY,GRAY,GRAY,NAVY,GRAY,sc_c,sc_c]; bolds=[True,False,False,False,True,False,True,True]
            for ci2,(v,al,co,bo) in enumerate(zip(vals,aligns,colors,bolds)):
                cell=tbl.cell(ri+1,ci2); cell_txt(cell,v,sz=7,bold=bo,color=co,align=al)
                if ci2==6: bar_fill(cell,p['real'],s)
                else: cell.fill.solid(); cell.fill.fore_color.rgb=bg

        # Guardar PPTX
        buf=io.BytesIO(); prs.save(buf); buf.seek(0)
        pptx_bytes=buf.read()

        # ── Modificar slide 3 con eventualidades del CSV ─────────────────────
        try:
            import zipfile as _zf, csv as _csv, io as _io2

            # Leer eventualidades_imc.csv del repo
            ev_url='https://raw.githubusercontent.com/DTI-Dashboard/Proyectos/main/eventualidades_imc.csv?t='+str(int(__import__('time').time()))
            try:
                with _ur.urlopen(ev_url) as _re2:
                    ev_content=_re2.read().decode('utf-8')
            except:
                ev_content=''

            ev_list=[]
            if ev_content.strip():
                reader=_csv.DictReader(_io2.StringIO(ev_content))
                for row in reader:
                    proy=(row.get('proyecto') or '').strip()
                    text=(row.get('eventualidad') or '').strip()
                    if proy and text:
                        ev_list.append((proy, text))

            if ev_list:
                # Leer el PPTX que acabamos de generar
                zip_in=_zf.ZipFile(_io2.BytesIO(pptx_bytes))
                s3_xml=zip_in.read('ppt/slides/slide3.xml').decode('utf-8')

                # Construir la tabla con la estructura exacta del template (compatible WPS)
                BORDER=('<a:lnL w="12700" cap="flat" cmpd="sng" algn="ctr"><a:solidFill><a:schemeClr val="tx1"/></a:solidFill><a:prstDash val="solid"/></a:lnL>'
                        '<a:lnR w="12700" cap="flat" cmpd="sng" algn="ctr"><a:solidFill><a:schemeClr val="tx1"/></a:solidFill><a:prstDash val="solid"/></a:lnR>'
                        '<a:lnT w="12700" cap="flat" cmpd="sng" algn="ctr"><a:solidFill><a:schemeClr val="tx1"/></a:solidFill><a:prstDash val="solid"/></a:lnT>'
                        '<a:lnB w="12700" cap="flat" cmpd="sng" algn="ctr"><a:solidFill><a:schemeClr val="tx1"/></a:solidFill><a:prstDash val="solid"/></a:lnB>')

                def ev_cell(txt, align, sz, bold=False):
                    b=' b="1"' if bold else ''
                    safe=txt.replace('&','&amp;').replace('<','&lt;').replace('>','&gt;').replace('"','&quot;')
                    return (f'<a:tc><a:txBody><a:bodyPr/><a:lstStyle/><a:p>'
                            f'<a:pPr marL="0" algn="{align}" defTabSz="457200" rtl="0" eaLnBrk="1" fontAlgn="b" latinLnBrk="0" hangingPunct="1"><a:buNone/></a:pPr>'
                            f'<a:r><a:rPr lang="es-MX" sz="{sz}"{b} u="none" strike="noStrike" kern="1200">'
                            f'<a:solidFill><a:schemeClr val="tx1"/></a:solidFill><a:effectLst/>'
                            f'<a:latin typeface="Aptos (cuerpo)"/></a:rPr>'
                            f'<a:t>{safe}</a:t></a:r></a:p></a:txBody>'
                            f'<a:tcPr marL="6244" marR="6244" marT="6244" marB="0" anchor="ctr">{BORDER}<a:noFill/></a:tcPr></a:tc>')

                import re as _re2
                tbl_hdr_m=_re2.search(r'<a:tbl>([\s\S]*?)<a:tr', s3_xml)
                orig_hdr=tbl_hdr_m.group(1) if tbl_hdr_m else '<a:tblPr/><a:tblGrid/>'

                hdr_row=('<a:tr h="520138">'
                         +ev_cell('No','ctr',1100,True)
                         +ev_cell('Proyecto','ctr',1100,True)
                         +ev_cell('Eventualidades','ctr',1100,True)
                         +'</a:tr>')
                data_rows=''
                for idx_ev,(proy,text) in enumerate(ev_list):
                    data_rows+=(f'<a:tr h="420000">'
                                +ev_cell(str(idx_ev+1),'ctr',1000)
                                +ev_cell(proy,'l',1000)
                                +ev_cell(text,'l',1000)
                                +'</a:tr>')

                new_tbl=f'<a:tbl>{orig_hdr}{hdr_row}{data_rows}</a:tbl>'
                s3_new=_re2.sub(r'<a:tbl>[\s\S]*?</a:tbl>', new_tbl, s3_xml)

                # Reempaquetar el PPTX
                buf2=_io2.BytesIO()
                zip_out=_zf.ZipFile(buf2,'w',_zf.ZIP_DEFLATED)
                for item in zip_in.infolist():
                    d=zip_in.read(item.filename)
                    if item.filename=='ppt/slides/slide3.xml':
                        d=s3_new.encode('utf-8')
                    zip_out.writestr(item,d)
                zip_out.close()
                pptx_bytes=buf2.getvalue()
                print(f"✅ Slide 3 con {len(ev_list)} eventualidades incluidas")
            else:
                print("ℹ️  Sin eventualidades en CSV — slide 3 sin modificar")
        except Exception as _es3:
            import traceback
            print(f"⚠️  Error en slide 3: {_es3}")
            print(traceback.format_exc())

        # Subir al repo via API
        api_url='https://api.github.com/repos/DTI-Dashboard/Proyectos/contents/cierre_base.pptx'
        gh_token=os.environ.get('GH_TOKEN','')
        headers_gh={'Authorization':f'token {gh_token}','Content-Type':'application/json'}
        try:
            req_get=_ur.Request(api_url,headers={'Authorization':f'token {gh_token}'})
            with _ur.urlopen(req_get) as rg: existing_sha=json.loads(rg.read())['sha']
        except: existing_sha=None
        payload={'message':'sync: actualizar cierre_base.pptx','content':base64.b64encode(pptx_bytes).decode()}
        if existing_sha: payload['sha']=existing_sha
        req_put=_ur.Request(api_url,data=json.dumps(payload).encode(),method='PUT',headers=headers_gh)
        with _ur.urlopen(req_put) as rp: json.loads(rp.read())
        print(f"✅ cierre_base.pptx generado: {len(datos)} proyectos, rango: {rango}")
    except Exception as _ep:
        import traceback
        print(f"⚠️  Error generando cierre_base.pptx: {_ep}")
        print(traceback.format_exc())

    # ── Timestamp ─────────────────────────────────────────────────
    ts_iso     = datetime.now(timezone.utc).strftime("%Y-%m-%dT%H:%M:%SZ")
    ts_display = datetime.now().strftime("%d/%m/%Y %H:%M")
    meta_tag   = f'<meta name="last-sync" content="{ts_iso}">'
    if 'name="last-sync"' in html:
        html = re.sub(r'<meta name="last-sync"[^>]*>', meta_tag, html)

    with open("index.html", "w", encoding="utf-8") as f:
        f.write(html)

    print(f"✅ index.html actualizado — {ts_display}")

    # Escribir éxito
    with open(_error_file, "w") as _f:
        _f.write("OK")
except Exception as _e:
    _err = _tb.format_exc()
    print("❌ EXCEPCIÓN:", _err)
    with open(_error_file, "w") as _f:
        _f.write(f"ERROR: {type(_e).__name__}: {_e}\n\n{_err}")
    _sys.exit(1)
