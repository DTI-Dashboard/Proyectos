"""
generar_reporte_semanal.py
──────────────────────────
Regenera BAJO DEMANDA la base de un reporte de cierre semanal:
  - Para un modulo individual (imc/dev/inf/ia): cierre_base[_mod].pptx con
    3 diapositivas utiles (portada, KPIs+tabla, eventualidades-placeholder).
  - Para 'resumen': cierre_base_resumen.pptx con 1 portada general + el
    bloque completo (KPIs+tabla y eventualidades-placeholder) de las 4
    direcciones, una tras otra (IMC, Desarrollo, Soporte, IA), + cierre.

La diapositiva de "eventualidades" NUNCA se llena aqui con datos reales:
eso lo arma el propio navegador en exportCierreSemanal() usando
localStorage (dti_eventualidades), que es la fuente de verdad real.

Se dispara desde el boton "Reporte Semanal (pptx)" de cada modulo (o de
Resumen), via el workflow .github/workflows/reporte_semanal.yml.

Requiere el secret GH_TOKEN (o PAT_TOKEN) con permisos de escritura sobre
el repo, para poder subir el pptx actualizado.
"""
import traceback as _tb
import sys as _sys

_error_file = "reporte_last_error.txt"

try:
    import os, json, re, base64, io, copy
    import urllib.request as _ur
    from datetime import datetime, date, timedelta

    GH_TOKEN = os.environ.get("GH_TOKEN", "") or os.environ.get("PAT_TOKEN", "")
    if not GH_TOKEN:
        raise Exception("Falta el secret GH_TOKEN (o PAT_TOKEN) en el workflow")

    # ── Determinar para que modulo se genera el reporte (por defecto IMC) ──
    MOD_DATA_VAR = {'imc':'IMC_DATA','dev':'DEV_DATA','inf':'INF_DATA','ia':'IA_DATA'}
    MOD_TITULO   = {'imc':'IMC — ÁREA DE IMPLEMENTACIÓN','dev':'DESARROLLO','inf':'SOPORTE E INFRAESTRUCTURA','ia':'INTELIGENCIA ARTIFICIAL'}
    MOD_ARCHIVO  = {'imc':'cierre_base.pptx','dev':'cierre_base_dev.pptx','inf':'cierre_base_inf.pptx','ia':'cierre_base_ia.pptx','resumen':'cierre_base_resumen.pptx'}
    MOD_SUBTITULO = {
        'imc': 'Dirección de Implementación y Mejora Continua',
        'dev': 'Dirección de Desarrollo',
        'inf': 'Dirección de Soporte e Infraestructura',
        'ia':  'Dirección de Inteligencia Artificial',
        'resumen': 'Resumen Ejecutivo — Todas las Direcciones',
    }
    ORDEN_RESUMEN = ['imc', 'dev', 'inf', 'ia']  # orden en que aparecen en el reporte combinado

    modulo = 'imc'
    try:
        req_mod = _ur.Request(
            "https://api.github.com/repos/DTI-Dashboard/Proyectos/contents/reporte_semanal_modulo.txt",
            headers={"Authorization": f"token {GH_TOKEN}"}
        )
        with _ur.urlopen(req_mod) as r:
            info_mod = json.loads(r.read())
        modulo_leido = base64.b64decode(info_mod["content"]).decode("utf-8").strip().lower()
        if modulo_leido in MOD_DATA_VAR or modulo_leido == 'resumen':
            modulo = modulo_leido
    except Exception:
        pass  # si el archivo no existe todavia, se queda en 'imc' (comportamiento original)

    archivo_salida = MOD_ARCHIVO[modulo]

    # ── Descargar el index.html actual del repo (ya trae los datos frescos del último sync) ──
    req_html = _ur.Request(
        "https://api.github.com/repos/DTI-Dashboard/Proyectos/contents/index.html",
        headers={"Authorization": f"token {GH_TOKEN}"}
    )
    with _ur.urlopen(req_html) as r:
        info = json.loads(r.read())
    with _ur.urlopen(info["download_url"]) as r:
        html = r.read().decode("utf-8")

    from pptx import Presentation
    from pptx.util import Inches as I, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.dml import MSO_PATTERN_TYPE
    from pptx.oxml.ns import qn
    import lxml.etree as etree

    def extraer_datos_modulo(html, mod_key):
        data_var = MOD_DATA_VAR[mod_key]
        idx = html.find(f'const {data_var} = [')
        if idx == -1:
            idx = html.find(f'const {data_var}=[')  # por si no tiene espacios (INF_DATA se guarda asi)
            i_pos = idx + len(f'const {data_var}=')
        else:
            i_pos = idx + len(f'const {data_var} = ')
        depth = 0
        block = ''
        for j in range(i_pos, i_pos + 500000):
            if html[j] == '[': depth += 1
            elif html[j] == ']':
                depth -= 1
                if depth == 0:
                    block = html[i_pos:j+1]
                    break
        matches = re.findall(
            r'\{tipo:"([^"]*)",nombre:"([^"]*)",area:"([^"]*)",inicio:"([^"]*)",fin:"([^"]*)",real:(\d+),esp:(\d+),sem:"([^"]*)"',
            block
        )
        return [{'tipo':m[0],'nombre':m[1],'area':m[2],'inicio':m[3],'fin':m[4],'real':int(m[5]),'esp':int(m[6]),'sem':m[7]} for m in matches]

    # Leer datos: un solo modulo, o los 4 si es 'resumen'
    if modulo == 'resumen':
        datos_por_mod = {mk: extraer_datos_modulo(html, mk) for mk in ORDEN_RESUMEN}
    else:
        datos_por_mod = {modulo: extraer_datos_modulo(html, modulo)}

    # Descargar template
    tpl_url = 'https://raw.githubusercontent.com/DTI-Dashboard/Proyectos/main/template_cierre.pptx'
    with _ur.urlopen(tpl_url) as _r:
        tpl_bytes = _r.read()

    hoy_d = date.today()
    dow = hoy_d.weekday()  # Lunes=0 ... Jueves=3 ... Domingo=6
    diff = dow - 3
    if diff > 3: diff -= 7
    if diff < -3: diff += 7
    jA = hoy_d - timedelta(days=diff)   # Jueves mas cercano a hoy (puede ser pasado o futuro)
    jP = jA - timedelta(days=7)
    MESES = ['Enero','Febrero','Marzo','Abril','Mayo','Junio','Julio','Agosto','Septiembre','Octubre','Noviembre','Diciembre']
    rango = f"{jP.day} de {MESES[jP.month-1]} al {jA.day} de {MESES[jA.month-1]}"

    prs = Presentation(io.BytesIO(tpl_bytes))
    NAVY=RGBColor(0x1A,0x23,0x40); GOLD=RGBColor(0xC9,0xA8,0x4C)
    WHITE=RGBColor(0xFF,0xFF,0xFF); GRAY=RGBColor(0x55,0x55,0x55)
    sem_col = {'Terminado':RGBColor(0x22,0xC5,0x5E),'A tiempo':RGBColor(0x22,0xC5,0x5E),'Atrasado':RGBColor(0xEF,0x44,0x44),'Stand By':RGBColor(0x11,0x11,0x11),'Riesgo de atraso':RGBColor(0xF5,0x9E,0x0B)}
    sem_hex = {'Terminado':'22C55E','A tiempo':'22C55E','Atrasado':'EF4444','Stand By':'111111','Riesgo de atraso':'F59E0B'}
    sem_txt = {'Terminado':'Terminado','A tiempo':'A tiempo','Atrasado':'Atrasado','Stand By':'Stand by','Riesgo de atraso':'Riesgo de atraso'}

    def sc(p):
        s=(p['sem'] or '').lower()
        if 'terminado' in s: return 'Terminado'
        if 'tiempo' in s: return 'A tiempo'
        if 'riesgo' in s: return 'Riesgo de atraso'
        if 'stand' in s: return 'Stand By'
        return 'Atrasado'

    def add_cell_text(slide, x, y, w, h, text, sz=7, bold=False, color=None, align=PP_ALIGN.CENTER):
        tb = slide.shapes.add_textbox(I(x+0.02), I(y), I(w-0.04), I(h))
        tf = tb.text_frame
        tf.word_wrap = False
        tf.margin_left=0; tf.margin_right=0; tf.margin_top=0; tf.margin_bottom=0
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.alignment = align
        r = p.add_run(); r.text = str(text)
        r.font.size = Pt(sz); r.font.bold = bold; r.font.name = 'Calibri'
        if color: r.font.color.rgb = color

    def add_row_bg(slide, x, y, w, h, color):
        rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, I(x), I(y), I(w), I(h))
        rect.fill.solid(); rect.fill.fore_color.rgb = color
        rect.line.fill.background()
        rect.shadow.inherit = False
        return rect

    def add_outer_shadow(shape, color_hex, blur_pt=6, dist_pt=1.5, alpha_pct=55, direction_deg=90):
        """Sombra exterior del MISMO color de la barra (no gris)."""
        spPr = shape._element.spPr
        for el in spPr.findall(qn('a:effectLst')):
            spPr.remove(el)
        effectLst = etree.SubElement(spPr, qn('a:effectLst'))
        outerShdw = etree.SubElement(effectLst, qn('a:outerShdw'))
        outerShdw.set('blurRad', str(int(blur_pt * 12700)))
        outerShdw.set('dist', str(int(dist_pt * 12700)))
        outerShdw.set('dir', str(int(direction_deg * 60000)))
        outerShdw.set('rotWithShape', '0')
        srgbClr = etree.SubElement(outerShdw, qn('a:srgbClr'))
        srgbClr.set('val', color_hex)
        alpha = etree.SubElement(srgbClr, qn('a:alpha'))
        alpha.set('val', str(int(alpha_pct * 1000)))

    def draw_bar(slide, x, y, w, h, pct, color_hex):
        """Barra tipo píldora: fondo con mini-cuadros claros del mismo color,
        relleno solido con sombra del mismo color, y % al final."""
        pct = max(0.0, min(100.0, float(pct or 0)))
        track = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, I(x), I(y), I(w), I(h))
        track.adjustments[0] = 0.5
        track.fill.patterned()
        track.fill.pattern = MSO_PATTERN_TYPE.SMALL_GRID
        track.fill.fore_color.rgb = RGBColor.from_string(color_hex)
        track.fill.back_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        track.line.fill.background()
        track.shadow.inherit = False

        if pct > 0:
            fill_w = max(h, w * (pct / 100.0))
            bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, I(x), I(y), I(fill_w), I(h))
            bar.adjustments[0] = 0.5
            bar.fill.solid()
            bar.fill.fore_color.rgb = RGBColor.from_string(color_hex)
            bar.line.fill.background()
            bar.shadow.inherit = False
            add_outer_shadow(bar, color_hex)

        txt = slide.shapes.add_textbox(I(x + w + 0.05), I(y - 0.035), I(0.5), I(h + 0.07))
        tf = txt.text_frame
        tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
        r = p.add_run(); r.text = f"{int(round(pct))}%"
        r.font.size = Pt(8); r.font.bold = True; r.font.name = 'Calibri'
        r.font.color.rgb = RGBColor.from_string(color_hex)

    def construir_slide_kpi(slide, datos, titulo_slide, rango):
        """Llena una diapositiva (vacia, layout de KPI+tabla) con los KPIs y
        la tabla de proyectos de un modulo especifico."""
        for shape in list(slide.shapes):
            try:
                if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE = 13
                    shape._element.getparent().remove(shape._element)
            except Exception:
                pass  # si el shape ya no tiene una relacion valida (p.ej. imagen de una slide duplicada), se ignora

        tot=len(datos); ter=sum(1 for p in datos if sc(p)=='Terminado')
        ati=sum(1 for p in datos if sc(p)=='A tiempo'); rie=sum(1 for p in datos if sc(p)=='Riesgo de atraso')
        atr=sum(1 for p in datos if sc(p)=='Atrasado')
        pcr=sum(1 for p in datos if p['real']>=90 and sc(p)!='Terminado')
        stb=sum(1 for p in datos if sc(p)=='Stand By')
        kpis=[('Total de proyectos',tot),('Terminado',ter),('A tiempo',ati),('En riesgo de atraso',rie),('Atrasados',atr),('Por cerrar',pcr),('Stand By',stb)]
        kw=1.345; kh=0.62; ky=0.62; kg=0.055; ksx=0.08
        for ii,(lbl,val) in enumerate(kpis):
            x=ksx+ii*(kw+kg)
            bx=slide.shapes.add_shape(1,I(x),I(ky),I(kw),I(kh))
            bx.fill.solid(); bx.fill.fore_color.rgb=WHITE
            bx.line.color.rgb=GOLD; bx.line.width=Pt(1.2)
            vt=slide.shapes.add_textbox(I(x),I(ky+0.03),I(kw),I(0.32))
            vp=vt.text_frame.paragraphs[0]; vp.alignment=PP_ALIGN.CENTER
            vr=vp.add_run(); vr.text=str(val); vr.font.size=Pt(20); vr.font.bold=True; vr.font.name='Calibri'; vr.font.color.rgb=NAVY
            lt=slide.shapes.add_textbox(I(x),I(ky+0.37),I(kw),I(0.24))
            lt.text_frame.word_wrap=True
            lp=lt.text_frame.paragraphs[0]; lp.alignment=PP_ALIGN.CENTER
            lr=lp.add_run(); lr.text=lbl; lr.font.size=Pt(6.5); lr.font.name='Calibri'; lr.font.color.rgb=GRAY
        tx=slide.shapes.add_textbox(I(1.75),I(0.04),I(6.5),I(0.34))
        p_tx=tx.text_frame.paragraphs[0]; p_tx.alignment=PP_ALIGN.CENTER
        r_tx=p_tx.add_run(); r_tx.text=titulo_slide
        r_tx.font.size=Pt(15); r_tx.font.bold=True; r_tx.font.name='Calibri'; r_tx.font.color.rgb=NAVY
        tx2=slide.shapes.add_textbox(I(0),I(0.38),I(10),I(0.18))
        p2=tx2.text_frame.paragraphs[0]; p2.alignment=PP_ALIGN.CENTER
        r2=p2.add_run(); r2.text=f'Cierre semanal del {rango}'
        r2.font.size=Pt(8.5); r2.font.name='Calibri'; r2.font.color.rgb=GRAY
        colW_in=[0.83,2.70,0.73,0.73,0.58,0.65,2.42,0.92]
        TBL_X=0.08; TBL_Y=1.30; HDR_H=0.22; ROW_H=0.205
        col_x=[TBL_X]
        for w in colW_in[:-1]:
            col_x.append(col_x[-1]+w)
        tblW=sum(colW_in)

        add_row_bg(slide, TBL_X, TBL_Y, tblW, HDR_H, GOLD)
        hdrs=['TIPO','PROYECTO','INICIO','FIN','% REAL','% ESPERADO','BARRA DE AVANCE','SEMÁFORO']
        for ci,h in enumerate(hdrs):
            add_cell_text(slide, col_x[ci], TBL_Y, colW_in[ci], HDR_H, h, sz=7, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

        bar_w=1.85
        bar_h=ROW_H*0.55
        for ri,p in enumerate(datos):
            s=sc(p); sc_c=sem_col.get(s,RGBColor(0x88,0x88,0x88))
            bg=RGBColor(0xF5,0xF7,0xFA) if ri%2==0 else WHITE
            row_y=TBL_Y+HDR_H+ri*ROW_H
            add_row_bg(slide, TBL_X, row_y, tblW, ROW_H, bg)
            nombre_corto=p['nombre'] if len(p['nombre'])<=48 else p['nombre'][:46]+'…'
            celdas=[
                (0,p['tipo'],PP_ALIGN.CENTER,NAVY,True),
                (1,nombre_corto,PP_ALIGN.LEFT,NAVY,False),
                (2,p['inicio'],PP_ALIGN.CENTER,GRAY,False),
                (3,p['fin'],PP_ALIGN.CENTER,GRAY,False),
                (4,f"{p['real']}%",PP_ALIGN.CENTER,NAVY,True),
                (5,f"{p['esp']}%",PP_ALIGN.CENTER,GRAY,False),
                (7,sem_txt.get(s,s),PP_ALIGN.CENTER,sc_c,True),
            ]
            for ci,v,al,co,bo in celdas:
                add_cell_text(slide, col_x[ci], row_y, colW_in[ci], ROW_H, v, sz=7, bold=bo, color=co, align=al)
            bar_row_y=row_y+(ROW_H-bar_h)/2
            draw_bar(slide, col_x[6]+0.05, bar_row_y, bar_w, bar_h, p['real'], sem_hex.get(s,'888888'))

    def duplicate_slide(prs, index):
        """Duplica una slide (por indice logico) y la agrega al FINAL. Devuelve la nueva slide.
        No copia shapes de tipo imagen: su relacion (rId) apunta al part original y no se
        puede duplicar de forma segura; el fondo real de estas slides viene del slide_layout,
        no de imagenes sueltas dentro de la propia slide."""
        source = prs.slides[index]
        dest = prs.slides.add_slide(source.slide_layout)
        for shp in list(dest.shapes):
            shp._element.getparent().remove(shp._element)
        for shp in source.shapes:
            try:
                if shp.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                    continue
            except Exception:
                pass
            newel = copy.deepcopy(shp._element)
            dest.shapes._spTree.append(newel)
        return dest

    def move_slide(prs, old_index, new_index):
        xml_slides = prs.slides._sldIdLst
        slides = list(xml_slides)
        xml_slides.remove(slides[old_index])
        xml_slides.insert(new_index, slides[old_index])

    def delete_slide(prs, index):
        xml_slides = prs.slides._sldIdLst
        slides = list(xml_slides)
        xml_slides.remove(slides[index])

    MOD_LABEL_CORTO = {'imc':'IMC','dev':'Desarrollo','inf':'Soporte e Infraestructura','ia':'IA'}

    if modulo != 'resumen':
        # ── Caso: un solo modulo (comportamiento original) ──
        subtitulo = MOD_SUBTITULO[modulo]
        titulo_slide = MOD_TITULO[modulo]
        datos = datos_por_mod[modulo]

        s1=prs.slides[0]
        for shape in s1.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if 'semanal del' in run.text:
                            run.text=f'Cierre semanal del {rango}'
                            shape.width=I(7.5); shape.left=I(2.1)
                            para.alignment=PP_ALIGN.RIGHT
                            run.font.name='Arial Nova Light'
                            run.font.size=Pt(13.9)
                            run.font.bold=True
                        elif 'Implementación y Mejora Continua' in run.text:
                            run.text = subtitulo

        s2=prs.slides[1]
        construir_slide_kpi(s2, datos, titulo_slide, rango)

        buf=io.BytesIO(); prs.save(buf); buf.seek(0)
        pptx_bytes_raw=buf.read()
        import zipfile as _zf, io as _io
        z_in=_zf.ZipFile(_io.BytesIO(pptx_bytes_raw))
        buf_out=_io.BytesIO()
        z_out=_zf.ZipFile(buf_out,'w',_zf.ZIP_DEFLATED)
        for item in z_in.infolist():
            d=z_in.read(item.filename)
            if item.filename=='ppt/slides/_rels/slide2.xml.rels':
                d=d.replace(b'slideLayouts/slideLayout12.xml',b'slideLayouts/slideLayout17.xml')
            elif item.filename=='ppt/slides/slide2.xml':
                s2_str=d.decode('utf-8')
                shapes=re.findall(r'<p:sp>[\s\S]*?</p:sp>',s2_str)
                for sh in shapes:
                    if 'ctrTitle' in sh:
                        s2_str=s2_str.replace(sh,'',1)
                        break
                d=s2_str.encode('utf-8')
            z_out.writestr(item,d)
        z_out.close()
        pptx_bytes=buf_out.getvalue()
        n_proyectos_msg = f"{len(datos)} proyectos"

    else:
        # ── Caso: 'resumen' — combinar las 4 direcciones en un solo pptx ──
        # Estructura del template: 0=portada, 1=kpi_template, 2=event_template, 3=cierre
        s1=prs.slides[0]
        for shape in s1.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if 'semanal del' in run.text:
                            run.text=f'Cierre semanal del {rango}'
                            shape.width=I(7.5); shape.left=I(2.1)
                            para.alignment=PP_ALIGN.RIGHT
                            run.font.name='Arial Nova Light'
                            run.font.size=Pt(13.9)
                            run.font.bold=True
                        elif 'Implementación y Mejora Continua' in run.text:
                            run.text = MOD_SUBTITULO['resumen']

        kpi_partnames = []    # nombres fisicos reales (ej. 'slide5.xml') de cada slide KPI, en orden
        event_partnames = []  # idem para las de eventualidades

        for mk in ORDEN_RESUMEN:
            kpi_slide = duplicate_slide(prs, 1)
            construir_slide_kpi(kpi_slide, datos_por_mod[mk], MOD_TITULO[mk], rango)
            kpi_partnames.append(kpi_slide.part.partname.split('/')[-1])

            event_slide = duplicate_slide(prs, 2)
            for shape in event_slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if 'Eventualidades de proyectos con atraso' in run.text:
                                run.text = f'Eventualidades de proyectos con atraso — {MOD_LABEL_CORTO[mk]}'
            event_partnames.append(event_slide.part.partname.split('/')[-1])

        # Mover "cierre" (indice logico 3) al final, y borrar las 2 slides plantilla originales
        move_slide(prs, 3, len(list(prs.slides)) - 1)
        delete_slide(prs, 2)  # event_template
        delete_slide(prs, 1)  # kpi_template

        buf=io.BytesIO(); prs.save(buf); buf.seek(0)
        pptx_bytes_raw=buf.read()
        import zipfile as _zf, io as _io
        z_in=_zf.ZipFile(_io.BytesIO(pptx_bytes_raw))
        buf_out=_io.BytesIO()
        z_out=_zf.ZipFile(buf_out,'w',_zf.ZIP_DEFLATED)
        for item in z_in.infolist():
            d=z_in.read(item.filename)
            fname=item.filename.split('/')[-1]
            if item.filename.startswith('ppt/slides/_rels/') and fname.replace('.rels','') in kpi_partnames:
                d=d.replace(b'slideLayouts/slideLayout12.xml',b'slideLayouts/slideLayout17.xml')
            elif item.filename.startswith('ppt/slides/') and fname in kpi_partnames:
                s_str=d.decode('utf-8')
                shapes=re.findall(r'<p:sp>[\s\S]*?</p:sp>',s_str)
                for sh in shapes:
                    if 'ctrTitle' in sh:
                        s_str=s_str.replace(sh,'',1)
                        break
                d=s_str.encode('utf-8')
            z_out.writestr(item,d)
        z_out.close()
        pptx_bytes=buf_out.getvalue()
        n_proyectos_msg = f"{sum(len(v) for v in datos_por_mod.values())} proyectos (4 direcciones) | kpi_slides={kpi_partnames} event_slides={event_partnames}"

    # Subir al repo via API
    api_url=f'https://api.github.com/repos/DTI-Dashboard/Proyectos/contents/{archivo_salida}'
    headers_gh={'Authorization':f'token {GH_TOKEN}','Content-Type':'application/json'}
    try:
        req_get=_ur.Request(api_url,headers={'Authorization':f'token {GH_TOKEN}'})
        with _ur.urlopen(req_get) as rg: existing_sha=json.loads(rg.read())['sha']
    except: existing_sha=None
    payload={'message':f'reporte: actualizar {archivo_salida} (bajo demanda, modulo {modulo})','content':base64.b64encode(pptx_bytes).decode()}
    if existing_sha: payload['sha']=existing_sha
    req_put=_ur.Request(api_url,data=json.dumps(payload).encode(),method='PUT',headers=headers_gh)
    with _ur.urlopen(req_put) as rp: json.loads(rp.read())
    print(f"✅ {archivo_salida} generado bajo demanda ({modulo}): {n_proyectos_msg}, rango: {rango}")

    with open(_error_file, "w") as _f:
        _f.write("OK")

except Exception as _e:
    _err = _tb.format_exc()
    print("❌ EXCEPCIÓN:", _err)
    with open(_error_file, "w") as _f:
        _f.write(f"ERROR: {type(_e).__name__}: {_e}\n\n{_err}")
    _sys.exit(1)
